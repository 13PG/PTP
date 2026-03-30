from __future__ import annotations

from datetime import datetime
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.services.topic_extractor import summarize_text

BACKGROUND = RGBColor(247, 243, 236)
TEXT_PRIMARY = RGBColor(35, 33, 28)
TEXT_MUTED = RGBColor(107, 98, 84)
ACCENT = RGBColor(201, 103, 52)
ACCENT_DARK = RGBColor(102, 53, 28)
CARD = RGBColor(255, 252, 246)


def build_ppt(
    topic: str,
    analyzed_documents: list[dict],
    global_summary: list[str],
    output_path: Path,
    progress_callback=None,
) -> list[str]:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    total_steps = 3 + len(analyzed_documents) + sum(1 for item in analyzed_documents if item["document"].images)
    completed_steps = 0
    warnings: list[str] = []

    _add_cover_slide(presentation, topic, analyzed_documents)
    completed_steps = _notify_build_progress(progress_callback, completed_steps, total_steps, "正在生成封面页。")
    _add_summary_slide(presentation, topic, global_summary)
    completed_steps = _notify_build_progress(progress_callback, completed_steps, total_steps, "正在生成主题摘要页。")

    for item in analyzed_documents:
        document = item["document"]
        try:
            _add_document_overview_slide(
                presentation,
                document,
                item["bullets"],
                item["score"],
                item.get("overview", ""),
            )
        except Exception as exc:
            warnings.append(f"{document.filename}: 信息页生成失败，已跳过。原因: {exc}")
        completed_steps = _notify_build_progress(
            progress_callback, completed_steps, total_steps, f"正在生成 {document.filename} 的信息页。"
        )
        if document.images:
            try:
                added_visual = _add_document_visual_slide(presentation, document, item["top_chunks"])
                if not added_visual:
                    warnings.append(f"{document.filename}: 未找到可用配图，已跳过图片页。")
            except Exception as exc:
                warnings.append(f"{document.filename}: 图片页生成失败，已跳过。原因: {exc}")
            completed_steps = _notify_build_progress(
                progress_callback, completed_steps, total_steps, f"正在生成 {document.filename} 的配图页。"
            )

    _add_finish_slide(presentation)
    _notify_build_progress(progress_callback, completed_steps, total_steps, "正在写出最终 PPT 文件。")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
    return warnings


def _add_cover_slide(presentation: Presentation, topic: str, analyzed_documents: list[dict]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_background(slide)
    _add_top_band(slide, "Paper Topic PPT")

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.4), Inches(1.6))
    title_frame = _prepare_text_frame(title_box.text_frame)
    title_frame.word_wrap = True
    paragraph = title_frame.paragraphs[0]
    paragraph.text = topic
    paragraph.font.name = "Microsoft YaHei"
    paragraph.font.size = Pt(24)
    paragraph.font.bold = True
    paragraph.font.color.rgb = TEXT_PRIMARY
    paragraph.alignment = PP_ALIGN.LEFT

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.4), Inches(1.2))
    sub_frame = _prepare_text_frame(sub_box.text_frame)
    sub = sub_frame.paragraphs[0]
    sub.text = f"共处理 {len(analyzed_documents)} 篇文档，已提取题目、作者、主题相关片段与论文配图"
    sub.font.name = "Microsoft YaHei"
    sub.font.size = Pt(14)
    sub.font.color.rgb = TEXT_MUTED
    sub.alignment = PP_ALIGN.LEFT

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.2), Inches(11.6), Inches(3.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = RGBColor(230, 221, 207)

    points = [
        "自动抽取与主题最相关的关键段落",
        "优先保留论文标题、作者与摘要线索",
        "将 PDF / Word 中的图片带入 PPT 页面",
        "适合手机端上传，多文件一次生成",
    ]
    box = slide.shapes.add_textbox(Inches(1.1), Inches(3.6), Inches(10.8), Inches(2.4))
    frame = _prepare_text_frame(box.text_frame)
    frame.word_wrap = True
    for index, point in enumerate(points):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = point
        paragraph.level = 0
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = TEXT_PRIMARY
        paragraph.space_after = Pt(8)
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = 1.08


def _add_summary_slide(presentation: Presentation, topic: str, global_summary: list[str]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_background(slide)
    _add_section_title(slide, "主题摘要", f"围绕“{topic}”筛出的跨文档关键内容")

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.9)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = RGBColor(230, 221, 207)

    box = slide.shapes.add_textbox(Inches(1.1), Inches(2.15), Inches(11.0), Inches(4.2))
    frame = _prepare_text_frame(box.text_frame)
    frame.word_wrap = True
    for index, bullet in enumerate(global_summary):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = TEXT_PRIMARY
        paragraph.space_after = Pt(8)
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = 1.14


def _add_document_overview_slide(
    presentation: Presentation, document, bullets: list[str], score: float, overview: str
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_background(slide)
    _add_section_title(slide, "论文信息", document.filename)

    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.7), Inches(1.3)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(255, 237, 223)
    header.line.color.rgb = RGBColor(232, 188, 155)

    title_box = slide.shapes.add_textbox(Inches(1.05), Inches(1.7), Inches(10.8), Inches(0.5))
    title_frame = _prepare_text_frame(title_box.text_frame)
    title = title_frame.paragraphs[0]
    title.text = document.title
    title.font.name = "Microsoft YaHei"
    title.font.size = Pt(22)
    title.font.bold = True
    title.font.color.rgb = ACCENT_DARK
    title.alignment = PP_ALIGN.LEFT

    author_box = slide.shapes.add_textbox(Inches(1.05), Inches(2.18), Inches(10.8), Inches(0.4))
    author_frame = _prepare_text_frame(author_box.text_frame)
    author = author_frame.paragraphs[0]
    author.text = f"作者: {document.authors}    主题相关度: {score}"
    author.font.name = "Microsoft YaHei"
    author.font.size = Pt(12)
    author.font.color.rgb = TEXT_MUTED
    author.alignment = PP_ALIGN.LEFT

    left_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.0), Inches(7.0), Inches(3.9)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CARD
    left_card.line.color.rgb = RGBColor(230, 221, 207)

    bullet_box = slide.shapes.add_textbox(Inches(1.05), Inches(3.25), Inches(6.45), Inches(3.3))
    bullet_frame = _prepare_text_frame(bullet_box.text_frame)
    bullet_frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = bullet_frame.paragraphs[0] if index == 0 else bullet_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = TEXT_PRIMARY
        paragraph.space_after = Pt(7)
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = 1.12

    right_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(3.0), Inches(4.4), Inches(3.9)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = RGBColor(255, 247, 239)
    right_card.line.color.rgb = RGBColor(232, 205, 180)

    info_lines = [
        f"文件类型: {document.file_type.upper()}",
        f"图片数量: {len(document.images)}",
        f"文本片段: {len(document.chunks)}",
        f"内容概述: {summarize_text(overview or document.abstract or '未识别摘要', 90)}",
    ]
    info_box = slide.shapes.add_textbox(Inches(8.4), Inches(3.32), Inches(3.7), Inches(3.1))
    info_frame = _prepare_text_frame(info_box.text_frame)
    info_frame.word_wrap = True
    for index, line in enumerate(info_lines):
        paragraph = info_frame.paragraphs[0] if index == 0 else info_frame.add_paragraph()
        paragraph.text = line
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(15)
        paragraph.font.color.rgb = TEXT_PRIMARY
        paragraph.space_after = Pt(10)
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = 1.1


def _add_document_visual_slide(presentation: Presentation, document, top_chunks: list) -> bool:
    valid_images = [image for image in document.images if image.path.exists()]
    if not valid_images:
        return False

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_background(slide)
    _add_section_title(slide, "论文配图", document.title)

    images = valid_images[:2]
    positions = [
        (Inches(0.85), Inches(1.65), Inches(5.4), Inches(4.2)),
        (Inches(6.95), Inches(1.65), Inches(5.4), Inches(4.2)),
    ]

    for image, (left, top, width, height) in zip(images, positions):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = RGBColor(230, 221, 207)
        _add_fitted_picture(
            slide,
            image.path,
            left + Inches(0.15),
            top + Inches(0.15),
            width - Inches(0.3),
            height - Inches(0.8),
        )

        caption_box = slide.shapes.add_textbox(
            left + Inches(0.18), top + height - Inches(0.55), width - Inches(0.36), Inches(0.4)
        )
        caption_frame = _prepare_text_frame(caption_box.text_frame, margin=0.02)
        caption = caption_frame.paragraphs[0]
        caption.text = f"{image.origin} | {summarize_text(image.caption or '论文配图', 42)}"
        caption.font.name = "Microsoft YaHei"
        caption.font.size = Pt(11)
        caption.font.color.rgb = TEXT_MUTED
        caption.alignment = PP_ALIGN.CENTER

    note_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(6.0), Inches(11.5), Inches(0.95)
    )
    note_card.fill.solid()
    note_card.fill.fore_color.rgb = RGBColor(255, 237, 223)
    note_card.line.color.rgb = RGBColor(232, 188, 155)

    note_box = slide.shapes.add_textbox(Inches(1.05), Inches(6.22), Inches(11.0), Inches(0.45))
    note_frame = _prepare_text_frame(note_box.text_frame)
    note = note_frame.paragraphs[0]
    top_note = top_chunks[0].text if top_chunks else document.abstract
    note.text = f"主题相关说明: {summarize_text(top_note or '未识别到更强相关片段。', 140)}"
    note.font.name = "Microsoft YaHei"
    note.font.size = Pt(14)
    note.font.color.rgb = ACCENT_DARK
    note.alignment = PP_ALIGN.LEFT
    return True


def _add_finish_slide(presentation: Presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_background(slide)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(11.3), Inches(5.1)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = CARD
    band.line.color.rgb = RGBColor(230, 221, 207)

    title_box = slide.shapes.add_textbox(Inches(1.35), Inches(2.05), Inches(10.5), Inches(0.7))
    title_frame = _prepare_text_frame(title_box.text_frame)
    title = title_frame.paragraphs[0]
    title.text = "PPT 生成完成"
    title.font.name = "Microsoft YaHei"
    title.font.size = Pt(28)
    title.font.bold = True
    title.font.color.rgb = TEXT_PRIMARY
    title.alignment = PP_ALIGN.CENTER

    desc_box = slide.shapes.add_textbox(Inches(1.65), Inches(3.1), Inches(9.9), Inches(1.5))
    desc_frame = _prepare_text_frame(desc_box.text_frame)
    desc = desc_frame.paragraphs[0]
    desc.text = "已按主题整合论文关键信息、标题、作者与图片。"
    desc.font.name = "Microsoft YaHei"
    desc.font.size = Pt(18)
    desc.font.color.rgb = TEXT_MUTED
    desc.alignment = PP_ALIGN.CENTER

    time_box = slide.shapes.add_textbox(Inches(1.65), Inches(4.1), Inches(9.9), Inches(0.7))
    time_frame = _prepare_text_frame(time_box.text_frame)
    timestamp = time_frame.paragraphs[0]
    timestamp.text = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    timestamp.font.name = "Microsoft YaHei"
    timestamp.font.size = Pt(14)
    timestamp.font.color.rgb = ACCENT
    timestamp.alignment = PP_ALIGN.CENTER


def _paint_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND


def _add_top_band(slide, text: str) -> None:
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.45), Inches(2.8), Inches(0.5)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.color.rgb = ACCENT
    box = slide.shapes.add_textbox(Inches(0.98), Inches(0.54), Inches(2.35), Inches(0.26))
    frame = _prepare_text_frame(box.text_frame, margin=0.01)
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def _add_section_title(slide, title: str, subtitle: str) -> None:
    _add_top_band(slide, title)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.92), Inches(11.2), Inches(0.55))
    title_frame = _prepare_text_frame(title_box.text_frame)
    p = title_frame.paragraphs[0]
    p.text = subtitle
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.LEFT


def _add_fitted_picture(slide, image_path: Path, left, top, width, height) -> None:
    with Image.open(image_path) as image:
        image_width, image_height = image.size

    width_value = int(width)
    height_value = int(height)
    scale = min(width_value / max(image_width, 1), height_value / max(image_height, 1))
    render_width = int(image_width * scale)
    render_height = int(image_height * scale)
    offset_left = left + int((width_value - render_width) / 2)
    offset_top = top + int((height_value - render_height) / 2)
    slide.shapes.add_picture(str(image_path), offset_left, offset_top, width=render_width, height=render_height)


def _prepare_text_frame(text_frame, margin: float = 0.05):
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)
    return text_frame


def _notify_build_progress(progress_callback, completed_steps: int, total_steps: int, message: str) -> int:
    next_step = completed_steps + 1
    if progress_callback is not None:
        progress = 78 + int((next_step / max(total_steps, 1)) * 20)
        progress_callback(min(progress, 98), message)
    return next_step
